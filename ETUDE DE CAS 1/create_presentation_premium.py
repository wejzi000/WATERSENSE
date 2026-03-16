"""
Génération de présentation PowerPoint PREMIUM avec images
Étude de Cas - Analyse de l'intégration des étudiants ISEM
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_premium_presentation():
    # Créer une présentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Couleurs du thème
    BLUE_DARK = RGBColor(31, 78, 121)
    BLUE_LIGHT = RGBColor(68, 114, 196)
    ORANGE = RGBColor(237, 125, 49)
    GREEN = RGBColor(112, 173, 71)
    RED = RGBColor(192, 0, 0)
    GRAY = RGBColor(89, 89, 89)
    WHITE = RGBColor(255, 255, 255)
    
    base_path = os.path.dirname(__file__)
    
    # ========== SLIDE 1: TITRE AVEC DESIGN ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond dégradé simulé avec formes
    bg1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = BLUE_DARK
    bg1.line.fill.background()
    
    # Bande décorative
    decoration = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), prs.slide_width, Inches(0.1))
    decoration.fill.solid()
    decoration.fill.fore_color.rgb = ORANGE
    decoration.line.fill.background()
    
    # Titre principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "ANALYSE DE L'INTÉGRATION\nDES ÉTUDIANTS"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-titre avec icône
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "📊 Étude de cas - ACM et Classification"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Info
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1.2))
    info_frame = info_box.text_frame
    info_frame.text = "Baromètre ISEM 2017\n625 étudiants • 92 variables • 4 étapes méthodologiques"
    p = info_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(8), Inches(7), Inches(1.5), Inches(0.3))
    date_frame = date_box.text_frame
    date_frame.text = "Février 2026"
    p = date_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.RIGHT
    
    # ========== SLIDE 2: SOMMAIRE VISUEL ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre avec bande
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = BLUE_DARK
    title_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "📋 PLAN DE PRÉSENTATION"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Contenu avec icônes
    items = [
        ("🎯", "1", "CONTEXTE ET OBJECTIFS", ORANGE),
        ("🧹", "2", "ÉTAPE 1 : Exploration et Nettoyage", BLUE_LIGHT),
        ("📐", "3", "ÉTAPE 2 : ACM Exploratoire", BLUE_LIGHT),
        ("🔍", "4", "ÉTAPE 3 : ACM Finale", BLUE_LIGHT),
        ("👥", "5", "ÉTAPE 4 : Classification", BLUE_LIGHT),
        ("✅", "6", "RÉSULTATS ET CONCLUSION", GREEN),
    ]
    
    y_pos = 1.8
    for icon, num, title, color in items:
        # Cercle numéro
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(y_pos), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        
        num_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos), Inches(0.5), Inches(0.5))
        num_frame = num_box.text_frame
        num_frame.text = num
        p = num_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        num_frame.vertical_anchor = 1  # Center
        
        # Texte
        text_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(7), Inches(0.5))
        text_frame = text_box.text_frame
        text_frame.text = f"{icon} {title}"
        p = text_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = GRAY
        text_frame.vertical_anchor = 1
        
        y_pos += 0.8
    
    # ========== SLIDE 3: CONTEXTE VISUEL ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = BLUE_DARK
    title_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "🎯 CONTEXTE DE L'ÉTUDE"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Boîtes avec ombres
    # Boîte DONNÉES
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(4), Inches(2.5))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(240, 248, 255)
    box1.line.color.rgb = BLUE_LIGHT
    box1.line.width = Pt(2)
    
    data_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(3.6), Inches(2.1))
    tf = data_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "📊 DONNÉES"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.space_after = Pt(10)
    
    infos = ["• 625 étudiants (L1, L2, L3)", "• Baromètre ISEM 2017", "• 92 variables initiales", "• Variables qualitatives"]
    for info in infos:
        p = tf.add_paragraph()
        p.text = info
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.space_after = Pt(8)
    
    # Boîte OBJECTIF
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(1.5), Inches(4), Inches(2.5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(255, 250, 240)
    box2.line.color.rgb = ORANGE
    box2.line.width = Pt(2)
    
    obj_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(3.6), Inches(2.1))
    tf = obj_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🎯 OBJECTIF"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Comprendre les facteurs d'intégration des étudiants et identifier des profils distincts"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.line_spacing = 1.3
    
    # Boîte MÉTHODOLOGIE
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(8.5), Inches(2.7))
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(240, 255, 240)
    box3.line.color.rgb = GREEN
    box3.line.width = Pt(2)
    
    method_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(8.1), Inches(2.3))
    tf = method_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🔬 MÉTHODOLOGIE EN 4 ÉTAPES"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "1️⃣ Nettoyage des données → 2️⃣ ACM Exploratoire → 3️⃣ ACM Finale & Tests → 4️⃣ Classification CAH"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    
    # ========== SLIDE 4: ÉTAPE 1 AVEC IMAGES ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre avec numéro
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = BLUE_DARK
    title_bg.line.fill.background()
    
    num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.2), Inches(0.6), Inches(0.6))
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = ORANGE
    num_circle.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(0.6), Inches(0.6))
    num_frame = num_box.text_frame
    num_frame.text = "1"
    p = num_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    num_frame.vertical_anchor = 1
    
    title_box = slide.shapes.add_textbox(Inches(1.3), Inches(0.25), Inches(8.2), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "EXPLORATION ET NETTOYAGE DES DONNÉES"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Images Étape 1
    try:
        # Image 1: Valeurs manquantes
        img1_path = os.path.join(base_path, "Etape_1_Exploration_Donnees", "1_valeurs_manquantes.png")
        if os.path.exists(img1_path):
            slide.shapes.add_picture(img1_path, Inches(0.5), Inches(1.2), width=Inches(3))
        
        # Image 2: Nb modalités
        img2_path = os.path.join(base_path, "Etape_1_Exploration_Donnees", "2_nb_modalites.png")
        if os.path.exists(img2_path):
            slide.shapes.add_picture(img2_path, Inches(3.7), Inches(1.2), width=Inches(3))
        
        # Image 3: Sélection
        img3_path = os.path.join(base_path, "Etape_1_Exploration_Donnees", "3_selection_variables.png")
        if os.path.exists(img3_path):
            slide.shapes.add_picture(img3_path, Inches(6.9), Inches(1.2), width=Inches(2.6))
    except Exception as e:
        print(f"Note: Certaines images de l'étape 1 n'ont pas pu être ajoutées: {e}")
    
    # Texte explicatif en bas
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🎯 Actions réalisées :"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.space_after = Pt(8)
    
    actions = [
        "✓ Analyse des valeurs manquantes et traitement",
        "✓ Recodage des variables (92 → ~15-20 variables pertinentes)",
        "✓ Contrôle qualité et validation des distributions"
    ]
    
    for action in actions:
        p = tf.add_paragraph()
        p.text = action
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY
        p.space_after = Pt(5)
    
    # ========== SLIDE 5: ÉTAPE 2 AVEC ÉBOULIS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = BLUE_DARK
    title_bg.line.fill.background()
    
    num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.2), Inches(0.6), Inches(0.6))
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = ORANGE
    num_circle.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(0.6), Inches(0.6))
    num_frame = num_box.text_frame
    num_frame.text = "2"
    p = num_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    num_frame.vertical_anchor = 1
    
    title_box = slide.shapes.add_textbox(Inches(1.3), Inches(0.25), Inches(8.2), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "ACM EXPLORATOIRE - Choix des axes"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Image éboulis (grande)
    try:
        img_eboulis = os.path.join(base_path, "Etape_2_ACM_Exploratoire", "eboulis_valeurs_propres.png")
        if os.path.exists(img_eboulis):
            slide.shapes.add_picture(img_eboulis, Inches(0.5), Inches(1.2), width=Inches(5.5))
    except Exception as e:
        print(f"Note: Image éboulis non ajoutée: {e}")
    
    # Boîte résultats à droite
    results_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.2), Inches(1.2), Inches(3.3), Inches(3.5))
    results_box.fill.solid()
    results_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    results_box.line.color.rgb = ORANGE
    results_box.line.width = Pt(3)
    
    text_box = slide.shapes.add_textbox(Inches(6.4), Inches(1.4), Inches(2.9), Inches(3.1))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "📊 INERTIES"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "Axe 1"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GRAY
    p.space_after = Pt(3)
    
    p = tf.add_paragraph()
    p.text = "75.82% ⭐⭐⭐"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Axes 1+2"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GRAY
    p.space_after = Pt(3)
    
    p = tf.add_paragraph()
    p.text = "86.37% ⭐⭐"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "✅ DÉCISION"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(5)
    
    p = tf.add_paragraph()
    p.text = "Retenir 2 AXES"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # Texte explicatif en bas
    expl_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2))
    tf = expl_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "💡 Interprétation : Avec 2 axes, on conserve 86% de l'information. La cassure nette dans l'éboulis confirme ce choix optimal."
    p.font.size = Pt(15)
    p.font.color.rgb = GRAY
    p.line_spacing = 1.3
    
    # ========== SLIDE 6: ÉTAPE 3 AVEC PLANS ACM ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = BLUE_DARK
    title_bg.line.fill.background()
    
    num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.2), Inches(0.6), Inches(0.6))
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = ORANGE
    num_circle.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(0.6), Inches(0.6))
    num_frame = num_box.text_frame
    num_frame.text = "3"
    p = num_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    num_frame.vertical_anchor = 1
    
    title_box = slide.shapes.add_textbox(Inches(1.3), Inches(0.25), Inches(8.2), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "ACM FINALE - Plans factoriels"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Plan des modalités
    try:
        img_modalites = os.path.join(base_path, "Etape_3_ACM_Finale", "plan_modalites_axes1_2.png")
        if os.path.exists(img_modalites):
            slide.shapes.add_picture(img_modalites, Inches(0.5), Inches(1.2), width=Inches(4.5))
    except Exception as e:
        print(f"Note: Plan modalités non ajouté: {e}")
    
    # Plan avec variable cible
    try:
        img_cible = os.path.join(base_path, "Etape_3_ACM_Finale", "plan_avec_variable_cible.png")
        if os.path.exists(img_cible):
            slide.shapes.add_picture(img_cible, Inches(5.2), Inches(1.2), width=Inches(4.3))
    except Exception as e:
        print(f"Note: Plan cible non ajouté: {e}")
    
    # Légendes
    legend1 = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(4.5), Inches(0.4))
    legend1_frame = legend1.text_frame
    legend1_frame.text = "Plan des modalités (catégories des variables)"
    p = legend1_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER
    
    legend2 = slide.shapes.add_textbox(Inches(5.2), Inches(4.8), Inches(4.3), Inches(0.4))
    legend2_frame = legend2.text_frame
    legend2_frame.text = "Plan avec variable cible (intégration)"
    p = legend2_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER
    
    # Texte explicatif
    expl_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    tf = expl_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "💡 Lecture : Les modalités proches ont des profils similaires. Les individus sont positionnés selon leurs caractéristiques."
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.line_spacing = 1.3
    
    # ========== SLIDE 7: TEST CHI² ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RED
    title_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "🔥 RÉSULTAT MAJEUR : TEST DU CHI²"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Grande boîte résultat significatif
    box_sig = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(8), Inches(2))
    box_sig.fill.solid()
    box_sig.fill.fore_color.rgb = RGBColor(240, 255, 240)
    box_sig.line.color.rgb = GREEN
    box_sig.line.width = Pt(4)
    
    sig_text = slide.shapes.add_textbox(Inches(1.3), Inches(1.8), Inches(7.4), Inches(1.4))
    tf = sig_text.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "✅ SEXE : χ² = 14.76, p = 0.002 (**)  →  SIGNIFICATIF"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Seul facteur discriminant statistiquement l'intégration"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Boîte résultats non significatifs
    box_nonsig = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4), Inches(8), Inches(2.5))
    box_nonsig.fill.solid()
    box_nonsig.fill.fore_color.rgb = RGBColor(255, 245, 245)
    box_nonsig.line.color.rgb = RED
    box_nonsig.line.width = Pt(3)
    
    nonsig_text = slide.shapes.add_textbox(Inches(1.3), Inches(4.3), Inches(7.4), Inches(1.9))
    tf = nonsig_text.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "❌ TOUTES LES AUTRES VARIABLES : NON SIGNIFICATIF (p > 0.05)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RED
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)
    
    variables_ns = "Logement • Âge • Région • Sport • Niveau d'études • Parcours • Nationalité • Situation familiale • Temps de trajet • Association • Études des parents"
    
    p = tf.add_paragraph()
    p.text = variables_ns
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.4
    
    # ========== SLIDE 8: ÉTAPE 4 - DENDROGRAMME ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = BLUE_DARK
    title_bg.line.fill.background()
    
    num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.2), Inches(0.6), Inches(0.6))
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = ORANGE
    num_circle.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(0.6), Inches(0.6))
    num_frame = num_box.text_frame
    num_frame.text = "4"
    p = num_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    num_frame.vertical_anchor = 1
    
    title_box = slide.shapes.add_textbox(Inches(1.3), Inches(0.25), Inches(8.2), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "CLASSIFICATION - Dendrogramme CAH"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Dendrogramme (grande image)
    try:
        img_dendro = os.path.join(base_path, "Etape_4_Classification", "dendrogramme_CAH.png")
        if os.path.exists(img_dendro):
            slide.shapes.add_picture(img_dendro, Inches(0.5), Inches(1.2), width=Inches(5.5))
    except Exception as e:
        print(f"Note: Dendrogramme non ajouté: {e}")
    
    # Silhouette scores
    try:
        img_silhouette = os.path.join(base_path, "Etape_4_Classification", "silhouette_scores.png")
        if os.path.exists(img_silhouette):
            slide.shapes.add_picture(img_silhouette, Inches(6.2), Inches(1.2), width=Inches(3.3))
    except Exception as e:
        print(f"Note: Silhouette non ajouté: {e}")
    
    # Boîte décision
    decision_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.2), Inches(9), Inches(1.8))
    decision_box.fill.solid()
    decision_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    decision_box.line.color.rgb = ORANGE
    decision_box.line.width = Pt(3)
    
    decision_text = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(8.4), Inches(1.2))
    tf = decision_text.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "✅ CHOIX : k = 4 CLASSES"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Critères : cassure dendrogramme + silhouette scores + interprétabilité + effectifs suffisants"
    p.font.size = Pt(15)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # ========== SLIDE 9: PLAN DES CLASSES ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = BLUE_DARK
    title_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "👥 LES 4 PROFILS D'INTÉGRATION"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Plan des classes (grande image)
    try:
        img_classes = os.path.join(base_path, "Etape_4_Classification", "plan_classes.png")
        if os.path.exists(img_classes):
            slide.shapes.add_picture(img_classes, Inches(0.5), Inches(1.2), width=Inches(6))
    except Exception as e:
        print(f"Note: Plan classes non ajouté: {e}")
    
    # Tableau des effectifs à droite
    classes_data = [
        ("1", "ISOLÉS", "78", "12.5%", RED),
        ("2", "PEU INTÉGRÉS", "425", "68.0%", ORANGE),
        ("3", "BIEN INTÉGRÉS", "65", "10.4%", BLUE_LIGHT),
        ("4", "TRÈS INTÉGRÉS", "57", "9.1%", GREEN),
    ]
    
    y_start = 1.5
    for num, nom, effectif, pct, color in classes_data:
        # Barre de classe
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(y_start), Inches(2.7), Inches(0.8))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        # Texte classe
        text_box = slide.shapes.add_textbox(Inches(6.9), Inches(y_start + 0.05), Inches(2.5), Inches(0.7))
        tf = text_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"Classe {num} - {nom}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(3)
        
        p = tf.add_paragraph()
        p.text = f"n = {effectif} ({pct})"
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        
        y_start += 1.1
    
    # ========== SLIDE 10: PROFILS DÉTAILLÉS AVEC ARBRE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = BLUE_DARK
    title_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "🌳 ARBRE DE DÉCISION EXPLICATIF"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Arbre de décision (grande image)
    try:
        img_arbre = os.path.join(base_path, "Etape_4_Classification", "arbre_decision_explicatif.png")
        if os.path.exists(img_arbre):
            slide.shapes.add_picture(img_arbre, Inches(0.5), Inches(1.2), height=Inches(5.8))
    except Exception as e:
        print(f"Note: Arbre de décision non ajouté: {e}")
    
    # Texte explicatif à droite
    expl_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    tf = expl_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "📖 INTERPRÉTATION"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "L'arbre de décision identifie les règles qui caractérisent chaque classe."
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.space_after = Pt(15)
    p.line_spacing = 1.3
    
    p = tf.add_paragraph()
    p.text = "🔍 Variables clés :"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_LIGHT
    p.space_after = Pt(8)
    
    key_vars = [
        "• Heures loisirs campus",
        "• Fréquence soirées",
        "• Mode de travail",
        "• Usage LILLIAD",
        "• Participation associations"
    ]
    
    for var in key_vars:
        p = tf.add_paragraph()
        p.text = var
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.space_after = Pt(6)
    
    # ========== SLIDE 11: CARACTÉRISATION 4 PROFILS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = GREEN
    title_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "📊 CARACTÉRISATION DES 4 PROFILS"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Classe 1 - ISOLÉS
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(4.5), Inches(2.8))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(255, 240, 240)
    box1.line.color.rgb = RED
    box1.line.width = Pt(3)
    
    text1 = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(4.1), Inches(2.4))
    tf = text1.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🔴 CLASSE 1 - ISOLÉS (12.5%)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RED
    p.space_after = Pt(8)
    
    traits1 = [
        "✓ Usage intensif campus (+16h)",
        "✓ Membres associations",
        "✓ Travaillent seuls",
        "✗ Vie sociale limitée",
        "",
        "→ Investis académiquement",
        "   mais isolés socialement"
    ]
    
    for trait in traits1:
        p = tf.add_paragraph()
        p.text = trait
        if trait.startswith("→"):
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RED
        else:
            p.font.size = Pt(12)
            p.font.color.rgb = GRAY
        p.space_after = Pt(4)
    
    # Classe 2 - PEU INTÉGRÉS
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.3), Inches(4.3), Inches(2.8))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(255, 248, 230)
    box2.line.color.rgb = ORANGE
    box2.line.width = Pt(3)
    
    text2 = slide.shapes.add_textbox(Inches(5.4), Inches(1.5), Inches(3.9), Inches(2.4))
    tf = text2.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🟠 CLASSE 2 - PEU INTÉGRÉS (68%)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(8)
    
    traits2 = [
        "✓ Intégration moyenne",
        "✓ Quelques liens sociaux",
        "✓ Présence modérée",
        "",
        "→ Profil MAJORITAIRE",
        "   Ni isolés ni hyper-intégrés",
        "   Adaptation \"standard\""
    ]
    
    for trait in traits2:
        p = tf.add_paragraph()
        p.text = trait
        if trait.startswith("→"):
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = ORANGE
        else:
            p.font.size = Pt(12)
            p.font.color.rgb = GRAY
        p.space_after = Pt(4)
    
    # Classe 3 - BIEN INTÉGRÉS
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(4.5), Inches(2.8))
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(230, 240, 255)
    box3.line.color.rgb = BLUE_LIGHT
    box3.line.width = Pt(3)
    
    text3 = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(4.1), Inches(2.4))
    tf = text3.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🔵 CLASSE 3 - BIEN INTÉGRÉS (10.4%)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_LIGHT
    p.space_after = Pt(8)
    
    traits3 = [
        "✓ Bonne participation sociale",
        "✓ Réseau développé",
        "✓ Équilibre études/social",
        "",
        "→ Adaptation réussie",
        "   Intégration harmonieuse"
    ]
    
    for trait in traits3:
        p = tf.add_paragraph()
        p.text = trait
        if trait.startswith("→"):
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = BLUE_LIGHT
        else:
            p.font.size = Pt(12)
            p.font.color.rgb = GRAY
        p.space_after = Pt(4)
    
    # Classe 4 - TRÈS INTÉGRÉS
    box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(4.3), Inches(4.3), Inches(2.8))
    box4.fill.solid()
    box4.fill.fore_color.rgb = RGBColor(240, 255, 240)
    box4.line.color.rgb = GREEN
    box4.line.width = Pt(3)
    
    text4 = slide.shapes.add_textbox(Inches(5.4), Inches(4.5), Inches(3.9), Inches(2.4))
    tf = text4.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🟢 CLASSE 4 - TRÈS INTÉGRÉS (9.1%)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(8)
    
    traits4 = [
        "✓ Forte implication",
        "✓ Leaders sociaux",
        "✓ Nombreuses relations",
        "✓ Actifs associations",
        "",
        "→ Hyper-actifs socialement",
        "   Ambassadeurs intégration"
    ]
    
    for trait in traits4:
        p = tf.add_paragraph()
        p.text = trait
        if trait.startswith("→"):
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = GREEN
        else:
            p.font.size = Pt(12)
            p.font.color.rgb = GRAY
        p.space_after = Pt(4)
    
    # ========== SLIDE 12: SYNTHÈSE RÉSULTATS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = BLUE_DARK
    title_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "✅ SYNTHÈSE DES RÉSULTATS"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Points clés en boîtes
    # Boîte 1
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(4.3), Inches(1.8))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(255, 245, 245)
    box1.line.color.rgb = RED
    box1.line.width = Pt(2)
    
    text1 = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(3.9), Inches(1.4))
    tf = text1.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🔥 RÉSULTAT SURPRENANT"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RED
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Seul le SEXE discrimine significativement l'intégration"
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY
    p.line_spacing = 1.3
    
    # Boîte 2
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.3), Inches(4.3), Inches(1.8))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(240, 255, 240)
    box2.line.color.rgb = GREEN
    box2.line.width = Pt(2)
    
    text2 = slide.shapes.add_textbox(Inches(5.4), Inches(1.5), Inches(3.9), Inches(1.4))
    tf = text2.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "✅ SEGMENTATION RÉUSSIE"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "4 profils distincts et interprétables identifiés"
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY
    p.line_spacing = 1.3
    
    # Boîte 3
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.3), Inches(4.3), Inches(1.8))
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(255, 250, 240)
    box3.line.color.rgb = ORANGE
    box3.line.width = Pt(2)
    
    text3 = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(3.9), Inches(1.4))
    tf = text3.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "⚠️ GROUPE À RISQUE"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "12.5% des étudiants sont isolés et nécessitent un accompagnement"
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY
    p.line_spacing = 1.3
    
    # Boîte 4
    box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(3.3), Inches(4.3), Inches(1.8))
    box4.fill.solid()
    box4.fill.fore_color.rgb = RGBColor(240, 248, 255)
    box4.line.color.rgb = BLUE_LIGHT
    box4.line.width = Pt(2)
    
    text4 = slide.shapes.add_textbox(Inches(5.4), Inches(3.5), Inches(3.9), Inches(1.4))
    tf = text4.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "💡 POTENTIEL D'AMÉLIORATION"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_LIGHT
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "68% en intégration moyenne peuvent progresser"
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY
    p.line_spacing = 1.3
    
    # Implications
    impl_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    tf = impl_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "💬 IMPLICATIONS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "→ Actions ciblées nécessaires pour les Isolés  •  → Approfondir le rôle du genre  •  → Variables comportementales > sociodémographiques"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.line_spacing = 1.3
    
    # ========== SLIDE 13: RECOMMANDATIONS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = ORANGE
    title_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "💡 RECOMMANDATIONS"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Colonne gauche - Actions
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🎯 ACTIONS PRIORITAIRES"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)
    
    actions = [
        ("CLASSE 1 (Isolés)", [
            "• Programme de parrainage",
            "• Groupes de travail guidés",
            "• Suivi individualisé",
            "• Événements sociaux ciblés"
        ]),
        ("CLASSE 2 (Peu intégrés)", [
            "• Activités collectives variées",
            "• Encourager associations",
            "• Espaces de socialisation",
            "• Ateliers compétences sociales"
        ])
    ]
    
    for classe, liste_actions in actions:
        p = tf.add_paragraph()
        p.text = classe
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.space_after = Pt(6)
        
        for action in liste_actions:
            p = tf.add_paragraph()
            p.text = action
            p.font.size = Pt(13)
            p.font.color.rgb = GRAY
            p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        p.text = ""
        p.space_after = Pt(10)
    
    # Colonne droite - Approfondissements
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🔬 APPROFONDISSEMENTS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(12)
    
    approfondissements = [
        "• Étude qualitative sur le rôle du genre",
        "• Suivi longitudinal (évolution temporelle)",
        "• Analyse des réseaux sociaux",
        "• Impact des dispositifs d'accompagnement",
        "• Comparaison inter-établissements"
    ]
    
    for appr in approfondissements:
        p = tf.add_paragraph()
        p.text = appr
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "📊 LIMITES"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED
    p.space_after = Pt(10)
    
    limites = [
        "• Données transversales",
        "• Biais de déclaration possible",
        "• Nécessite validation qualitative"
    ]
    
    for limite in limites:
        p = tf.add_paragraph()
        p.text = limite
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.space_after = Pt(8)
    
    # ========== SLIDE 14: CONCLUSION ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond coloré
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BLUE_DARK
    background.line.fill.background()
    
    # Bandes décoratives
    deco1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), prs.slide_width, Inches(0.05))
    deco1.fill.solid()
    deco1.fill.fore_color.rgb = ORANGE
    deco1.line.fill.background()
    
    deco2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5), prs.slide_width, Inches(0.05))
    deco2.fill.solid()
    deco2.fill.fore_color.rgb = GREEN
    deco2.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = "CONCLUSION"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Points clés
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(7), Inches(2))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    conclusions = [
        "✓ Méthodologie rigoureuse en 4 étapes",
        "✓ Segmentation claire en 4 profils d'intégration",
        "✓ Identification d'un groupe à risque (12.5%)",
        "✓ Rôle significatif du sexe dans l'intégration",
        "✓ Base solide pour actions ciblées"
    ]
    
    for conclusion in conclusions:
        p = tf.add_paragraph()
        p.text = conclusion
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
    
    # Message final
    final_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    final_frame = final_box.text_frame
    final_frame.text = "Une analyse statistique complète au service de l'amélioration\nde l'intégration étudiante"
    p = final_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.4
    
    # ========== SLIDE 15: QUESTIONS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond coloré
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = ORANGE
    background.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "QUESTIONS ?"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-texte
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Merci de votre attention"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Contact
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    contact_frame = contact_box.text_frame
    contact_frame.text = "📧 Février 2026 • Baromètre ISEM 2017"
    p = contact_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Sauvegarder
    output_path = os.path.join(base_path, "PRESENTATION_PREMIUM_ETUDE_CAS_INTEGRATION.pptx")
    prs.save(output_path)
    print(f"✅ Présentation PREMIUM créée avec succès !")
    print(f"📍 Fichier: {output_path}")
    print(f"📊 {len(prs.slides)} slides générées avec images intégrées")
    return output_path

if __name__ == "__main__":
    create_premium_presentation()
