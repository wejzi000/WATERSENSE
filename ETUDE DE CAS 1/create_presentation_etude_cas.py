"""
Génération automatique de la présentation PowerPoint
Étude de Cas - Analyse de l'intégration des étudiants ISEM
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    # Créer une présentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Couleurs du thème
    BLUE_DARK = RGBColor(31, 78, 121)
    BLUE_LIGHT = RGBColor(68, 114, 196)
    ORANGE = RGBColor(237, 125, 49)
    GREEN = RGBColor(112, 173, 71)
    RED = RGBColor(192, 0, 0)
    GRAY = RGBColor(89, 89, 89)
    
    # ========== SLIDE 1: TITRE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Fond coloré
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLUE_DARK
    
    # Titre principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "ANALYSE DE L'INTÉGRATION DES ÉTUDIANTS"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Étude de cas - ACM et Classification"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Info
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
    info_frame = info_box.text_frame
    info_frame.text = "Baromètre ISEM 2017 • 625 étudiants • 4 étapes méthodologiques"
    p = info_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(8), Inches(7), Inches(1.5), Inches(0.3))
    date_frame = date_box.text_frame
    date_frame.text = "Février 2026"
    p = date_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.RIGHT
    
    # ========== SLIDE 2: SOMMAIRE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "PLAN DE PRÉSENTATION"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    # Contenu
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    items = [
        ("1", "CONTEXTE ET OBJECTIFS", "Population, problématique, méthodologie"),
        ("2", "ÉTAPE 1 : Exploration et Nettoyage", "92 variables → sélection et recodage"),
        ("3", "ÉTAPE 2 : ACM Exploratoire", "Dimensionnalité et choix des axes"),
        ("4", "ÉTAPE 3 : ACM Finale", "Interprétation et tests statistiques"),
        ("5", "ÉTAPE 4 : Classification", "Segmentation en 4 profils d'intégration"),
        ("6", "RÉSULTATS CLÉS ET CONCLUSION", "Synthèse et recommandations"),
    ]
    
    for i, (num, title, desc) in enumerate(items):
        p = tf.add_paragraph()
        p.text = f"{num}. {title}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = BLUE_LIGHT
        p.space_after = Pt(5)
        
        p = tf.add_paragraph()
        p.text = f"   {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.space_after = Pt(15)
    
    # ========== SLIDE 3: CONTEXTE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "CONTEXTE DE L'ÉTUDE"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    # Colonne gauche
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "📊 DONNÉES"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(10)
    
    infos = [
        "• 625 étudiants (L1, L2, L3)",
        "• Baromètre ISEM 2017",
        "• 92 variables initiales",
        "• Variables qualitatives",
        "",
        "🎯 OBJECTIF",
        "Comprendre les facteurs",
        "d'intégration des étudiants",
        "et identifier des profils",
        "distincts"
    ]
    
    for info in infos:
        if info.startswith("🎯"):
            p = tf.add_paragraph()
            p.text = info
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = ORANGE
            p.space_after = Pt(10)
            p.space_before = Pt(15)
        else:
            p = tf.add_paragraph()
            p.text = info
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY
            p.space_after = Pt(8)
    
    # Colonne droite
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.3), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🔬 MÉTHODOLOGIE"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(10)
    
    methods = [
        "1️⃣ Nettoyage des données",
        "   Gestion valeurs manquantes",
        "   Recodage variables",
        "",
        "2️⃣ ACM Exploratoire",
        "   Réduction dimensionnalité",
        "   Sélection axes principaux",
        "",
        "3️⃣ ACM Finale",
        "   Interprétation plan factoriel",
        "   Tests Chi²",
        "",
        "4️⃣ Classification CAH",
        "   Segmentation 4 profils",
        "   Caractérisation groupes"
    ]
    
    for method in methods:
        p = tf.add_paragraph()
        p.text = method
        if method.startswith(("1️⃣", "2️⃣", "3️⃣", "4️⃣")):
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = BLUE_LIGHT
        else:
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
        p.space_after = Pt(5)
    
    # ========== SLIDE 4: ÉTAPE 1 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Numéro
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1), Inches(0.8))
    num_frame = num_box.text_frame
    num_frame.text = "01"
    p = num_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(1.8), Inches(0.5), Inches(7.7), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "EXPLORATION ET NETTOYAGE DES DONNÉES"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    # Contenu
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5.4))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🎯 OBJECTIF"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_LIGHT
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Préparer des données propres et exploitables pour l'Analyse des Correspondances Multiples"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.space_after = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "⚙️ ACTIONS RÉALISÉES"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_LIGHT
    p.space_after = Pt(10)
    
    actions = [
        ("✓", "Analyse des valeurs manquantes", "Identification et traitement des données incomplètes"),
        ("✓", "Recodage des variables", "Regroupement des modalités (ex: âge → tranches d'âge)"),
        ("✓", "Sélection des variables", "92 variables → ~15-20 variables pertinentes"),
        ("✓", "Contrôle qualité", "Vérification cohérence et distributions"),
    ]
    
    for check, action, detail in actions:
        p = tf.add_paragraph()
        p.text = f"{check} {action}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p.space_after = Pt(5)
        
        p = tf.add_paragraph()
        p.text = f"   → {detail}"
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "📊 RÉSULTAT : Données structurées prêtes pour l'analyse statistique"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_before = Pt(15)
    
    # ========== SLIDE 5: ÉTAPE 2 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Numéro
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1), Inches(0.8))
    num_frame = num_box.text_frame
    num_frame.text = "02"
    p = num_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(1.8), Inches(0.5), Inches(7.7), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "ACM EXPLORATOIRE"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    # Contenu gauche
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(5.4))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🎯 OBJECTIF"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_LIGHT
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Réduire la dimensionnalité et sélectionner les axes synthétiques"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "📐 QU'EST-CE QU'UNE ACM ?"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_LIGHT
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Analyse des Correspondances Multiples = ACP pour variables qualitatives"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "• Résume l'information de 15+ variables en 2-3 axes"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    
    p = tf.add_paragraph()
    p.text = "• Identifie les associations entre modalités"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    
    p = tf.add_paragraph()
    p.text = "• Permet visualisation 2D"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    
    # Contenu droite - Résultats
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5.4))
    tf = right_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "📊 RÉSULTATS INERTIES"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(10)
    
    # Tableau résultats
    results = [
        ("Axe 1", "75.82%", "⭐⭐⭐"),
        ("Axes 1+2", "86.37%", "⭐⭐"),
        ("Axes 1+2+3", "93.54%", "⭐"),
    ]
    
    for axe, inertie, stars in results:
        p = tf.add_paragraph()
        p.text = f"{axe}: {inertie} {stars}"
        p.font.size = Pt(16)
        p.font.bold = True
        if axe == "Axes 1+2":
            p.font.color.rgb = ORANGE
        else:
            p.font.color.rgb = GRAY
        p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "✅ DÉCISION"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Retenir 2 AXES"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Justification :"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GRAY
    p.space_after = Pt(5)
    
    justifs = [
        "• 86% de l'information conservée",
        "• Cassure nette dans l'éboulis",
        "• Lisibilité (plan 2D)",
        "• Gain faible si axe 3 (+7%)"
    ]
    
    for j in justifs:
        p = tf.add_paragraph()
        p.text = j
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.space_after = Pt(5)
    
    # ========== SLIDE 6: ÉTAPE 3 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Numéro
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1), Inches(0.8))
    num_frame = num_box.text_frame
    num_frame.text = "03"
    p = num_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(1.8), Inches(0.5), Inches(7.7), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "ACM FINALE ET INTERPRÉTATION"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    # Contenu
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5.4))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🎯 OBJECTIF"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_LIGHT
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Interpréter le plan factoriel et identifier les variables discriminantes"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.space_after = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "📊 ANALYSES RÉALISÉES"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_LIGHT
    p.space_after = Pt(10)
    
    analyses = [
        "• Plan des modalités : Visualisation des catégories sur axes 1-2",
        "• Plan des individus : Position de chaque étudiant (625 points)",
        "• Matrice de corrélations : Liens entre variables",
        "• Tests du Chi² : Significativité des associations avec l'intégration"
    ]
    
    for analyse in analyses:
        p = tf.add_paragraph()
        p.text = analyse
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "🔥 RÉSULTAT MAJEUR : TEST DU CHI²"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED
    p.space_after = Pt(15)
    p.space_before = Pt(20)
    
    # Tableau résultats Chi²
    p = tf.add_paragraph()
    p.text = "✅ SEXE : χ² = 14.76, p = 0.002 → SIGNIFICATIF (**)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "   → Seul facteur discriminant statistiquement l'intégration"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "❌ Toutes les autres variables : p > 0.05 → NON SIGNIFICATIF"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RED
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "   → Logement, âge, région, sport, niveau, parcours, etc."
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    
    # ========== SLIDE 7: ÉTAPE 4 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Numéro
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1), Inches(0.8))
    num_frame = num_box.text_frame
    num_frame.text = "04"
    p = num_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(1.8), Inches(0.5), Inches(7.7), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "CLASSIFICATION - 4 PROFILS"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    # Contenu gauche
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(5.4))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🎯 OBJECTIF"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_LIGHT
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Créer des groupes homogènes d'étudiants selon leur profil d'intégration"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "⚙️ MÉTHODOLOGIE"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_LIGHT
    p.space_after = Pt(8)
    
    methods_class = [
        "1. Classification Ascendante Hiérarchique (CAH)",
        "2. Dendrogramme pour choix k",
        "3. Silhouette scores (qualité)",
        "4. Valeurs-tests (profiling)",
        "5. Arbre de décision explicatif"
    ]
    
    for method in methods_class:
        p = tf.add_paragraph()
        p.text = method
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "✅ CHOIX : k = 4 classes"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(5)
    
    p = tf.add_paragraph()
    p.text = "Basé sur critères statistiques et interprétabilité"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    
    # Contenu droite - 4 Classes
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5.4))
    tf = right_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "📊 LES 4 PROFILS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(12)
    
    classes = [
        ("Classe 1", "ISOLÉS", "78 (12.5%)", RED),
        ("Classe 2", "PEU INTÉGRÉS", "425 (68.0%)", ORANGE),
        ("Classe 3", "BIEN INTÉGRÉS", "65 (10.4%)", BLUE_LIGHT),
        ("Classe 4", "TRÈS INTÉGRÉS", "57 (9.1%)", GREEN),
    ]
    
    for num, nom, effectif, color in classes:
        p = tf.add_paragraph()
        p.text = num
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = GRAY
        p.space_after = Pt(3)
        
        p = tf.add_paragraph()
        p.text = nom
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(3)
        
        p = tf.add_paragraph()
        p.text = f"n = {effectif}"
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(5)
    
    p = tf.add_paragraph()
    p.text = "💡 À RETENIR"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "• 68% sont peu intégrés (profil majoritaire)"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    
    p = tf.add_paragraph()
    p.text = "• 12.5% en difficulté (isolés)"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    
    p = tf.add_paragraph()
    p.text = "• 19.5% bien/très intégrés"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    
    # ========== SLIDE 8: PROFILS DÉTAILLÉS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "CARACTÉRISATION DES 4 PROFILS"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    # Classe 1
    box1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.6))
    tf = box1.text_frame
    tf.word_wrap = True
    
    # Fond
    fill = box1.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 240, 240)
    
    p = tf.paragraphs[0]
    p.text = "CLASSE 1 - ISOLÉS (12.5%)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RED
    p.space_after = Pt(8)
    
    traits = [
        "✓ Usage intensif campus (+16h loisirs)",
        "✓ Membres associations",
        "✓ Travaillent seuls",
        "✗ Vie sociale limitée",
        "→ Profil : investis académiquement",
        "   mais isolés socialement"
    ]
    
    for trait in traits:
        p = tf.add_paragraph()
        p.text = trait
        p.font.size = Pt(11)
        if trait.startswith("→"):
            p.font.bold = True
            p.font.color.rgb = RED
        else:
            p.font.color.rgb = GRAY
        p.space_after = Pt(4)
    
    # Classe 2
    box2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(2.6))
    tf = box2.text_frame
    tf.word_wrap = True
    
    # Fond
    fill = box2.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 248, 230)
    
    p = tf.paragraphs[0]
    p.text = "CLASSE 2 - PEU INTÉGRÉS (68%)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(8)
    
    traits2 = [
        "✓ Intégration moyenne",
        "✓ Quelques liens sociaux",
        "✓ Présence modérée campus",
        "→ Profil : MAJORITAIRE",
        "   Ni isolés ni hyper-intégrés",
        "   Adaptation \"standard\""
    ]
    
    for trait in traits2:
        p = tf.add_paragraph()
        p.text = trait
        p.font.size = Pt(11)
        if trait.startswith("→"):
            p.font.bold = True
            p.font.color.rgb = ORANGE
        else:
            p.font.color.rgb = GRAY
        p.space_after = Pt(4)
    
    # Classe 3
    box3 = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(4.5), Inches(2.6))
    tf = box3.text_frame
    tf.word_wrap = True
    
    # Fond
    fill = box3.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 240, 255)
    
    p = tf.paragraphs[0]
    p.text = "CLASSE 3 - BIEN INTÉGRÉS (10.4%)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_LIGHT
    p.space_after = Pt(8)
    
    traits3 = [
        "✓ Bonne participation sociale",
        "✓ Réseau développé",
        "✓ Équilibre études/social",
        "→ Profil : Adaptation réussie",
        "   Intégration harmonieuse",
        "   Liens solides"
    ]
    
    for trait in traits3:
        p = tf.add_paragraph()
        p.text = trait
        p.font.size = Pt(11)
        if trait.startswith("→"):
            p.font.bold = True
            p.font.color.rgb = BLUE_LIGHT
        else:
            p.font.color.rgb = GRAY
        p.space_after = Pt(4)
    
    # Classe 4
    box4 = slide.shapes.add_textbox(Inches(5.2), Inches(4.3), Inches(4.3), Inches(2.6))
    tf = box4.text_frame
    tf.word_wrap = True
    
    # Fond
    fill = box4.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 255, 240)
    
    p = tf.paragraphs[0]
    p.text = "CLASSE 4 - TRÈS INTÉGRÉS (9.1%)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(8)
    
    traits4 = [
        "✓ Forte implication",
        "✓ Leaders sociaux",
        "✓ Nombreuses relations",
        "✓ Actifs associations",
        "→ Profil : Hyper-actifs",
        "   Ambassadeurs de l'intégration"
    ]
    
    for trait in traits4:
        p = tf.add_paragraph()
        p.text = trait
        p.font.size = Pt(11)
        if trait.startswith("→"):
            p.font.bold = True
            p.font.color.rgb = GREEN
        else:
            p.font.color.rgb = GRAY
        p.space_after = Pt(4)
    
    # ========== SLIDE 9: RÉSULTATS CLÉS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "RÉSULTATS CLÉS ET ENSEIGNEMENTS"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    # Contenu
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5.4))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🔥 RÉSULTAT SURPRENANT"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RED
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Seul le SEXE discrimine significativement l'intégration (p=0.002)"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.space_after = Pt(5)
    
    p = tf.add_paragraph()
    p.text = "Les variables classiques (logement, âge, région, sport, niveau...) ne sont PAS significatives"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.space_after = Pt(25)
    
    p = tf.add_paragraph()
    p.text = "✅ SEGMENTATION RÉUSSIE"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(10)
    
    segments = [
        "4 profils distincts et interprétables",
        "Séparation basée sur comportements sociaux et usage campus",
        "Identification d'un groupe à risque (Isolés, 12.5%)",
        "Majorité en intégration moyenne (68%) → potentiel d'amélioration"
    ]
    
    for seg in segments:
        p = tf.add_paragraph()
        p.text = f"• {seg}"
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "💡 IMPLICATIONS"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(10)
    p.space_before = Pt(15)
    
    implications = [
        "Nécessité d'actions ciblées pour les Isolés",
        "Potentiel d'amélioration pour 68% des étudiants",
        "Question du rôle du genre dans l'intégration à approfondir",
        "Variables comportementales > variables sociodémographiques"
    ]
    
    for impl in implications:
        p = tf.add_paragraph()
        p.text = f"→ {impl}"
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY
        p.space_after = Pt(8)
    
    # ========== SLIDE 10: RECOMMANDATIONS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "RECOMMANDATIONS"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    
    # Colonne gauche
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(5.4))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🎯 ACTIONS PRIORITAIRES"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(10)
    
    actions_prio = [
        ("CLASSE 1 (Isolés)", [
            "Programme de parrainage",
            "Groupes de travail guidés",
            "Suivi individualisé",
            "Événements sociaux ciblés"
        ]),
        ("CLASSE 2 (Peu intégrés)", [
            "Activités collectives",
            "Encourager associations",
            "Espaces de socialisation"
        ])
    ]
    
    for classe, actions in actions_prio:
        p = tf.add_paragraph()
        p.text = classe
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = BLUE_LIGHT
        p.space_after = Pt(6)
        
        for action in actions:
            p = tf.add_paragraph()
            p.text = f"• {action}"
            p.font.size = Pt(13)
            p.font.color.rgb = GRAY
            p.space_after = Pt(5)
        
        p = tf.add_paragraph()
        p.text = ""
        p.space_after = Pt(8)
    
    # Colonne droite
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5.4))
    tf = right_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🔬 APPROFONDISSEMENTS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(10)
    
    approfondissements = [
        "Étude qualitative pour comprendre le rôle du genre",
        "Suivi longitudinal (évolution des profils dans le temps)",
        "Analyse des réseaux sociaux",
        "Impact des dispositifs d'accompagnement"
    ]
    
    for appr in approfondissements:
        p = tf.add_paragraph()
        p.text = f"• {appr}"
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "📊 LIMITES DE L'ÉTUDE"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED
    p.space_after = Pt(10)
    
    limites = [
        "Données transversales (pas de causalité)",
        "Biais de déclaration possible",
        "Variables peut-être mal mesurées",
        "Nécessite validation qualitative"
    ]
    
    for limite in limites:
        p = tf.add_paragraph()
        p.text = f"• {limite}"
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.space_after = Pt(8)
    
    # ========== SLIDE 11: CONCLUSION ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond coloré
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLUE_DARK
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "CONCLUSION"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Contenu
    content_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    conclusions = [
        "✓ Méthodologie rigoureuse en 4 étapes",
        "✓ Segmentation claire en 4 profils d'intégration",
        "✓ Identification d'un groupe à risque (12.5%)",
        "✓ Résultat surprenant sur le rôle du sexe",
        "✓ Base solide pour actions ciblées"
    ]
    
    for conclusion in conclusions:
        p = tf.add_paragraph()
        p.text = conclusion
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(15)
    
    # Message final
    final_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8))
    final_frame = final_box.text_frame
    final_frame.text = "Une analyse statistique complète au service de l'amélioration de l'intégration étudiante"
    p = final_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    # ========== SLIDE 12: QUESTIONS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond coloré
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ORANGE
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "QUESTIONS ?"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-texte
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Merci de votre attention"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Sauvegarder
    output_path = os.path.join(os.path.dirname(__file__), "PRESENTATION_ETUDE_CAS_INTEGRATION_ETUDIANTS.pptx")
    prs.save(output_path)
    print(f"✅ Présentation créée avec succès : {output_path}")
    print(f"📊 {len(prs.slides)} slides générées")
    return output_path

if __name__ == "__main__":
    create_presentation()
